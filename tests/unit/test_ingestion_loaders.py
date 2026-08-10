"""Unit tests for per-format document loaders."""

import csv
import io

import docx
import openpyxl
import pymupdf
from pptx import Presentation

from app.ingestion.base import LoadedPage
from app.ingestion.csv_loader import CsvLoader
from app.ingestion.docx_loader import DocxLoader
from app.ingestion.excel_loader import ExcelLoader
from app.ingestion.pdf_loader import PdfLoader
from app.ingestion.pptx_loader import PptxLoader
from app.ingestion.text_loader import TextLoader


def test_pdf_loader_extracts_text_per_page() -> None:
    """`PdfLoader` should return one page per PDF page, with matching page numbers."""
    document = pymupdf.open()
    for text in ("Page one", "Page two"):
        page = document.new_page()
        page.insert_text((72, 72), text)
    data = document.tobytes()
    document.close()

    loaded = PdfLoader().load(data, "sample.pdf")

    assert loaded.source == "sample.pdf"
    assert loaded.content_type == "application/pdf"
    assert len(loaded.pages) == 2
    assert loaded.pages[0].page_number == 1
    assert "Page one" in loaded.pages[0].text
    assert loaded.pages[1].page_number == 2
    assert "Page two" in loaded.pages[1].text


def test_docx_loader_extracts_paragraph_text() -> None:
    """`DocxLoader` should join all paragraph text into a single page."""
    word_document = docx.Document()
    word_document.add_paragraph("Hello world")
    word_document.add_paragraph("Second paragraph")
    buffer = io.BytesIO()
    word_document.save(buffer)

    loaded = DocxLoader().load(buffer.getvalue(), "sample.docx")

    assert len(loaded.pages) == 1
    assert loaded.pages[0].page_number is None
    assert "Hello world" in loaded.pages[0].text
    assert "Second paragraph" in loaded.pages[0].text


def test_text_loader_decodes_plain_text() -> None:
    """`TextLoader` should decode UTF-8 bytes as-is for `.txt` files."""
    loaded = TextLoader(content_type="text/plain").load(b"plain text content", "sample.txt")

    assert loaded.content_type == "text/plain"
    assert loaded.pages == [LoadedPage(text="plain text content", page_number=None)]


def test_text_loader_decodes_markdown() -> None:
    """`TextLoader` should decode UTF-8 bytes as-is for `.md` files."""
    loaded = TextLoader(content_type="text/markdown").load(b"# Heading\nBody", "sample.md")

    assert loaded.content_type == "text/markdown"
    assert loaded.pages[0].text == "# Heading\nBody"


def test_csv_loader_joins_rows_as_text() -> None:
    """`CsvLoader` should render each CSV row as a comma-joined line."""
    buffer = io.StringIO()
    writer = csv.writer(buffer)
    writer.writerow(["name", "age"])
    writer.writerow(["Alice", "30"])
    data = buffer.getvalue().encode("utf-8")

    loaded = CsvLoader().load(data, "sample.csv")

    assert loaded.content_type == "text/csv"
    assert loaded.pages[0].text == "name, age\nAlice, 30"


def test_excel_loader_extracts_one_page_per_sheet() -> None:
    """`ExcelLoader` should produce one page per worksheet, in workbook order."""
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Sheet1"
    sheet.append(["a", "b"])
    workbook.create_sheet("Sheet2").append(["x", "y"])
    buffer = io.BytesIO()
    workbook.save(buffer)

    loaded = ExcelLoader().load(buffer.getvalue(), "sample.xlsx")

    assert len(loaded.pages) == 2
    assert loaded.pages[0].page_number == 1
    assert loaded.pages[0].text == "a, b"
    assert loaded.pages[1].page_number == 2
    assert loaded.pages[1].text == "x, y"


def test_pptx_loader_extracts_one_page_per_slide() -> None:
    """`PptxLoader` should produce one page per slide with all text-frame content."""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Title 1"
    slide.placeholders[1].text_frame.text = "Body text"
    buffer = io.BytesIO()
    presentation.save(buffer)

    loaded = PptxLoader().load(buffer.getvalue(), "sample.pptx")

    assert len(loaded.pages) == 1
    assert loaded.pages[0].page_number == 1
    assert "Title 1" in loaded.pages[0].text
    assert "Body text" in loaded.pages[0].text
