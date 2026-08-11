"""PowerPoint document loader using python-pptx. Each slide is treated as one page."""

import io

from pptx import Presentation

from app.ingestion.base import LoadedDocument, LoadedPage


class PptxLoader:
    """Loads text from PowerPoint (.pptx) presentations, one page per slide."""

    def load(self, data: bytes, filename: str) -> LoadedDocument:
        """Extract text from every text-bearing shape on each slide."""
        presentation = Presentation(io.BytesIO(data))
        pages = []
        for index, slide in enumerate(presentation.slides):
            texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
            pages.append(LoadedPage(text="\n".join(texts), page_number=index + 1))
        content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        return LoadedDocument(source=filename, content_type=content_type, pages=pages)
